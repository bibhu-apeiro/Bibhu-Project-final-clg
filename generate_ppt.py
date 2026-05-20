from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "XAI_Credit_Risk_Presentation.pptx")

PRIMARY = RGBColor(0x4F, 0x46, 0xE5)
PRIMARY_DARK = RGBColor(0x1E, 0x1B, 0x4B)
PRIMARY_LIGHT = RGBColor(0xEE, 0xF2, 0xFF)
ACCENT_CYAN = RGBColor(0x06, 0xB6, 0xD4)
ACCENT_CYAN_LIGHT = RGBColor(0xEC, 0xFE, 0xFF)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_GREEN_LIGHT = RGBColor(0xEC, 0xFD, 0xF5)
ACCENT_RED = RGBColor(0xF4, 0x3F, 0x5E)
ACCENT_RED_LIGHT = RGBColor(0xFF, 0xF1, 0xF2)
ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_AMBER_LIGHT = RGBColor(0xFF, 0xFB, 0xEB)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT_PURPLE_LIGHT = RGBColor(0xFA, 0xF5, 0xFF)
TEXT_DARK = RGBColor(0x1E, 0x29, 0x3B)
TEXT_MEDIUM = RGBColor(0x47, 0x55, 0x69)
TEXT_LIGHT = RGBColor(0x94, 0xA3, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)


def bg(slide, color=BG_LIGHT):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def rrect(slide, l, t, w, h, color, border_c=None, border_w=Pt(0)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if border_c:
        s.line.color.rgb = border_c
        s.line.width = border_w
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def oval(slide, l, t, w, h, color, border_c=None, border_w=Pt(0)):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if border_c:
        s.line.color.rgb = border_c
        s.line.width = border_w
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def chevron(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def diamond(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def hexagon(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def arrow_right(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def shape_text(shape, text, size=12, bold=False, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri"):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align


def shape_multitext(shape, lines, size=12, bold=False, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri"):
    tf = shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font
        p.alignment = align


def txbox(slide, l, t, w, h, text, size=14, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return tb


def bullets(slide, l, t, w, h, items, size=12, color=TEXT_DARK, spacing=Pt(6)):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        run = p.add_run()
        run.text = "\u2022  " + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def header_bar(slide, title, subtitle=None):
    rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.15), PRIMARY_DARK)
    txbox(slide, Inches(0.8), Inches(0.18), Inches(11), Inches(0.5), title, size=24, bold=True, color=WHITE)
    if subtitle:
        txbox(slide, Inches(0.8), Inches(0.62), Inches(11), Inches(0.35), subtitle, size=12, color=ACCENT_CYAN)
    rect(slide, Inches(0), Inches(1.15), Inches(13.33), Inches(0.06), ACCENT_CYAN)


def slide_num(slide, n, total=10):
    txbox(slide, Inches(12.2), Inches(7.1), Inches(1), Inches(0.3), f"{n} / {total}", size=9, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)


def icon_circle(slide, x, y, size, bg_color, symbol, symbol_color=WHITE, symbol_size=16):
    c = oval(slide, x, y, size, size, bg_color)
    shape_text(c, symbol, size=symbol_size, bold=True, color=symbol_color)
    return c


prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# =============================================================================
# SLIDE 1: TITLE SLIDE (Dark, dramatic)
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)

rect(s, Inches(0), Inches(0), Inches(13.33), Inches(0.12), ACCENT_CYAN)
rect(s, Inches(0), Inches(7.38), Inches(13.33), Inches(0.12), ACCENT_CYAN)

rrect(s, Inches(1.2), Inches(0.8), Inches(10.93), Inches(5.6), PRIMARY, ACCENT_CYAN, Pt(3))

txbox(s, Inches(1.5), Inches(1.1), Inches(10.33), Inches(0.5), "8th SEMESTER  RESEARCH  PROJECT", size=13, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

txbox(s, Inches(1.5), Inches(1.8), Inches(10.33), Inches(1.4), "Explainable AI Framework\nfor Credit Risk Prediction", size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(s, Inches(4.5), Inches(3.35), Inches(4.33), Inches(0.05), ACCENT_CYAN)

txbox(s, Inches(1.5), Inches(3.65), Inches(10.33), Inches(0.5), "Using Ensemble Learning & SHAP Explainability", size=17, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

for i, icon in enumerate(["\U0001F9E0", "\U0001F4CA", "\U0001F6E0", "\u26A1"]):
    icon_circle(s, Inches(3.5) + Inches(i * 1.6), Inches(4.4), Inches(0.55), ACCENT_CYAN, icon, PRIMARY_DARK, 18)

txbox(s, Inches(1.5), Inches(5.1), Inches(10.33), Inches(0.5), "Bibhu Prasad Sahoo  &  Soumyashree", size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txbox(s, Inches(1.5), Inches(5.7), Inches(10.33), Inches(0.3), "Under the Guidance of Research Faculty", size=11, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

tech_tags = ["Python", "Scikit-Learn", "SHAP", "FastAPI", "React", "Docker"]
tag_w = Inches(1.5)
total_w = len(tech_tags) * 1.7 - 0.2
start_x = (13.33 - total_w) / 2
for i, tag in enumerate(tech_tags):
    x = Inches(start_x + i * 1.7)
    rrect(s, x, Inches(6.15), tag_w, Inches(0.35), PRIMARY_DARK, ACCENT_CYAN, Pt(1))
    txbox(s, x, Inches(6.15), tag_w, Inches(0.35), tag, size=9, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 2: PROBLEM STATEMENT (Visual infographic style)
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header_bar(s, "The Problem", "Why do banks need Explainable AI?")
slide_num(s, 2)

txbox(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.5), "Imagine a bank rejects your loan. You ask \"Why?\" and the answer is... \"The computer said NO.\"", size=16, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

problems = [
    ("\U0001F512", "Black Box AI", "AI says YES or NO, but\ncannot explain WHY.", "Banks don't know\nwhat influenced the\ndecision.", ACCENT_RED, ACCENT_RED_LIGHT),
    ("\u26A0", "Legal Risk", "Laws like GDPR & RBI rules\nsay: \"You MUST explain\nloan decisions.\"", "Banks can face heavy\nfines for unexplained\nrejections.", ACCENT_AMBER, ACCENT_AMBER_LIGHT),
    ("\U0001F614", "No Trust", "If people don't understand\nthe reason, they lose\ntrust in the bank.", "Applicants can't fix\ntheir profile if they\ndon't know what's wrong.", PRIMARY, PRIMARY_LIGHT),
    ("\U0001F52A", "Accuracy vs Clarity", "Complex AI = more accurate\nbut harder to explain.", "Simple AI = easy to explain\nbut less accurate.", ACCENT_PURPLE, ACCENT_PURPLE_LIGHT),
]

for i, (icon, title, line1, line2, accent, light_bg) in enumerate(problems):
    x = Inches(0.4) + (i * Inches(3.15))
    rrect(s, x, Inches(2.2), Inches(2.95), Inches(4.2), WHITE, BORDER, Pt(1))
    rect(s, x, Inches(2.2), Inches(2.95), Inches(0.08), accent)
    icon_circle(s, x + Inches(0.95), Inches(2.5), Inches(0.9), accent, icon, WHITE, 24)
    txbox(s, x + Inches(0.15), Inches(3.55), Inches(2.65), Inches(0.4), title, size=14, bold=True, color=accent, align=PP_ALIGN.CENTER)
    txbox(s, x + Inches(0.15), Inches(3.95), Inches(2.65), Inches(1.0), line1, size=10, color=TEXT_MEDIUM, align=PP_ALIGN.CENTER)
    rrect(s, x + Inches(0.15), Inches(5.1), Inches(2.65), Inches(1.1), light_bg, accent, Pt(1))
    txbox(s, x + Inches(0.25), Inches(5.2), Inches(2.45), Inches(0.9), line2, size=10, bold=True, color=accent, align=PP_ALIGN.CENTER)

rrect(s, Inches(1.5), Inches(6.6), Inches(10.33), Inches(0.55), PRIMARY, PRIMARY, Pt(0))
txbox(s, Inches(1.5), Inches(6.62), Inches(10.33), Inches(0.55), "SOLUTION:  Build a system that is BOTH accurate AND explainable", size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 3: OBJECTIVES (Roadmap style with connected steps)
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header_bar(s, "Our Goals", "What we set out to build")
slide_num(s, 3)

objectives = [
    ("\U0001F3AF", "Build Smart AI", "Train 3 different AI models\nthat vote together like a jury", ACCENT_GREEN),
    ("\U0001F50D", "Make It Transparent", "Use SHAP to show exactly\nwhich factors drove the decision", ACCENT_CYAN),
    ("\U0001F4BB", "Build a Dashboard", "Create a web app where users\ncan test and explore results live", PRIMARY),
    ("\U0001F916", "Add AI Summaries", "Auto-generate plain English\nexplanations for every decision", ACCENT_AMBER),
    ("\U0001F4C4", "Export PDF Reports", "One-click professional reports\nwith full analysis for audits", ACCENT_PURPLE),
]

for i, (icon, title, desc, accent) in enumerate(objectives):
    x = Inches(0.3) + (i * Inches(2.55))
    rrect(s, x, Inches(1.5), Inches(2.35), Inches(4.0), WHITE, BORDER, Pt(1))
    rect(s, x, Inches(1.5), Inches(2.35), Inches(0.06), accent)
    icon_circle(s, x + Inches(0.72), Inches(1.75), Inches(0.8), accent, icon, WHITE, 22)
    txbox(s, x + Inches(0.1), Inches(2.7), Inches(2.15), Inches(0.4), title, size=13, bold=True, color=accent, align=PP_ALIGN.CENTER)
    txbox(s, x + Inches(0.1), Inches(3.2), Inches(2.15), Inches(1.2), desc, size=10, color=TEXT_MEDIUM, align=PP_ALIGN.CENTER)

    if i < len(objectives) - 1:
        arrow_right(s, x + Inches(2.35), Inches(3.0), Inches(0.2), Inches(0.3), TEXT_LIGHT)

step_labels = ["STEP 1", "STEP 2", "STEP 3", "STEP 4", "STEP 5"]
for i, label in enumerate(step_labels):
    x = Inches(0.3) + (i * Inches(2.55))
    rrect(s, x + Inches(0.7), Inches(5.7), Inches(0.95), Inches(0.35), PRIMARY_DARK, None, Pt(0))
    txbox(s, x + Inches(0.7), Inches(5.72), Inches(0.95), Inches(0.3), label, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(s, Inches(0.8), Inches(6.15), Inches(11.73), Inches(0.03), PRIMARY)

txbox(s, Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.8),
      "End Result:  A complete system that predicts loan risk accurately AND explains every decision in plain English", size=14, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 4: TECHNOLOGY STACK (Icon grid)
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header_bar(s, "Tech Stack", "The tools we used to build this system")
slide_num(s, 4)

sections = [
    ("\U0001F527", "Backend (Server Side)", PRIMARY, [
        ("\U0001F40D", "Python", "Core programming language"),
        ("\u26A1", "FastAPI", "Fast web API framework"),
        ("\U0001F9E0", "Scikit-Learn", "Machine learning library"),
        ("\U0001F4CA", "SHAP", "Explainability engine"),
        ("\U0001F4C8", "Pandas / NumPy", "Data handling"),
        ("\U0001F4C4", "ReportLab", "PDF generation"),
    ]),
    ("\U0001F5A5", "Frontend (User Interface)", ACCENT_CYAN, [
        ("\u269B", "React", "Interactive UI framework"),
        ("\u26A1", "Vite", "Ultra-fast build tool"),
        ("\U0001F3A8", "Tailwind CSS", "Beautiful styling"),
        ("\U0001F4C9", "Recharts", "Charts & graphs"),
        ("\U0001F4E1", "Axios", "API communication"),
        ("\U0001F310", "Nginx", "Web server"),
    ]),
    ("\U0001F433", "Deployment", ACCENT_GREEN, [
        ("\U0001F4E6", "Docker", "App packaging"),
        ("\U0001F517", "Docker Compose", "Running multiple containers"),
        ("\U0001F504", "Hot Reload", "Live code updates"),
        ("\U0001F517", "CORS", "Secure cross-origin"),
        ("\U0001F4CD", "Ports", "Frontend: 3000, API: 8000"),
        ("\U0001F512", "SSL Ready", "Secure deployment"),
    ]),
    ("\U0001F4CB", "Input Data (9 Features)", ACCENT_AMBER, [
        ("\U0001F464", "Age", "Applicant's age"),
        ("\U0001F4B0", "Income", "Yearly earnings (\u20b9)"),
        ("\U0001F4B5", "Loan Amount", "Money requested"),
        ("\U0001F4C8", "Credit Score", "Risk rating (300-850)"),
        ("\U0001F3E0", "Experience", "Years working"),
        ("\U0001F393", "Education", "Qualification level"),
    ]),
]

for si, (icon, title, accent, items) in enumerate(sections):
    x = Inches(0.3) + (si * Inches(3.2))
    rrect(s, x, Inches(1.35), Inches(3.0), Inches(5.7), WHITE, BORDER, Pt(1))
    rect(s, x, Inches(1.35), Inches(3.0), Inches(0.55), accent)
    txbox(s, x + Inches(0.1), Inches(1.38), Inches(2.8), Inches(0.5), f"{icon}  {title}", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for ii, (iicon, name, desc) in enumerate(items):
        iy = Inches(2.1) + (ii * Inches(0.78))
        rrect(s, x + Inches(0.1), iy, Inches(2.8), Inches(0.65), BG_LIGHT, BORDER, Pt(1))
        txbox(s, x + Inches(0.15), iy + Inches(0.0), Inches(2.7), Inches(0.3), f"{iicon}  {name}", size=10, bold=True, color=accent, align=PP_ALIGN.LEFT)
        txbox(s, x + Inches(0.4), iy + Inches(0.3), Inches(2.45), Inches(0.3), desc, size=8, color=TEXT_MEDIUM, align=PP_ALIGN.LEFT)

# =============================================================================
# SLIDE 5: SYSTEM ARCHITECTURE (Visual flow diagram)
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header_bar(s, "How It Works", "The complete journey from loan application to decision")
slide_num(s, 5)

flow_items = [
    ("\U0001F464", "Applicant", "Fills in loan\ndetails", PRIMARY_DARK),
    ("\u2192", "", "", None),
    ("\U0001F4BB", "Dashboard", "React web\ninterface", ACCENT_CYAN),
    ("\u2192", "", "", None),
    ("\u2699", "API Server", "FastAPI\nprocesses data", PRIMARY),
    ("\u2192", "", "", None),
    ("\U0001F9E0", "3 AI Agents", "Each votes\nApprove/Reject", ACCENT_GREEN),
    ("\u2192", "", "", None),
    ("\U0001F4CA", "SHAP Engine", "Explains WHY\nthe decision", ACCENT_PURPLE),
    ("\u2192", "", "", None),
    ("\U0001F4AC", "AI Summary", "Plain English\nreport", ACCENT_AMBER),
]

xi = Inches(0.15)
for icon, title, desc, accent in flow_items:
    if icon == "\u2192":
        arrow_right(s, xi, Inches(1.55), Inches(0.35), Inches(0.5), TEXT_LIGHT)
        xi += Inches(0.35)
    else:
        rrect(s, xi, Inches(1.35), Inches(1.85), Inches(1.5), WHITE, accent, Pt(2))
        icon_circle(s, xi + Inches(0.55), Inches(1.45), Inches(0.65), accent, icon, WHITE, 18)
        txbox(s, xi + Inches(0.05), Inches(2.15), Inches(1.75), Inches(0.3), title, size=10, bold=True, color=accent, align=PP_ALIGN.CENTER)
        txbox(s, xi + Inches(0.05), Inches(2.45), Inches(1.75), Inches(0.4), desc, size=8, color=TEXT_MEDIUM, align=PP_ALIGN.CENTER)
        xi += Inches(1.85)

txbox(s, Inches(0.5), Inches(3.2), Inches(12.33), Inches(0.4), "Step-by-Step: What happens when you click \"Analyze\"", size=15, bold=True, color=PRIMARY_DARK, align=PP_ALIGN.CENTER)

steps = [
    ("\u2460", "You enter details", "Age, Income, Loan Amount,\nCredit Score, Education, etc.", ACCENT_CYAN),
    ("\u2461", "Data is cleaned", "Missing values filled, text\nconverted to numbers, scaled", PRIMARY),
    ("\u2462", "3 AI models vote", "Random Forest, Logistic Regression\n& Neural Network each predict", ACCENT_GREEN),
    ("\u2463", "Votes are combined", "Average of all 3 gives the\nfinal approval probability", ACCENT_PURPLE),
    ("\u2464", "SHAP explains it", "Shows which factors helped\n(+) or hurt (-) your chances", ACCENT_AMBER),
    ("\u2465", "You get the answer", "Clear verdict, visual charts,\nactionable tips, and PDF report", ACCENT_RED),
]

for i, (num, title, desc, accent) in enumerate(steps):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + (col * Inches(4.15))
    y = Inches(3.75) + (row * Inches(1.7))
    rrect(s, x, y, Inches(3.95), Inches(1.5), WHITE, accent, Pt(2))
    icon_circle(s, x + Inches(0.15), y + Inches(0.35), Inches(0.7), accent, num, WHITE, 22)
    txbox(s, x + Inches(1.0), y + Inches(0.15), Inches(2.8), Inches(0.35), title, size=13, bold=True, color=accent)
    txbox(s, x + Inches(1.0), y + Inches(0.55), Inches(2.8), Inches(0.8), desc, size=10, color=TEXT_MEDIUM)

# =============================================================================
# SLIDE 6: ML PIPELINE (Visual with 3 agent cards)
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header_bar(s, "The 3 AI Agents", "How our ensemble system makes decisions")
slide_num(s, 6)

agents = [
    ("\U0001F333", "Agent 1: Random Forest", ACCENT_GREEN, ACCENT_GREEN_LIGHT, [
        "Creates 100 decision trees",
        "Each tree votes on the application",
        "Great at finding hidden patterns",
        "Captures complex relationships",
        "Lead agent for SHAP analysis",
    ], "Like asking 100 experts\nand taking the majority vote"),
    ("\U0001F4CF", "Agent 2: Logistic Regression", ACCENT_CYAN, ACCENT_CYAN_LIGHT, [
        "A simple, transparent model",
        "Creates a linear decision boundary",
        "Easy to understand and explain",
        "Provides reliable baseline",
        "Good for regulatory reporting",
    ], "Like a simple checklist:\nmore points = higher chance"),
    ("\U0001F9E0", "Agent 3: Neural Network", ACCENT_PURPLE, ACCENT_PURPLE_LIGHT, [
        "2 hidden layers (64 & 32 neurons)",
        "Learns very complex patterns",
        "Finds subtle connections",
        "Complements the other agents",
        "Catches edge cases",
    ], "Like a brain that learns\nfrom thousands of examples"),
]

for i, (icon, title, accent, light, items, analogy) in enumerate(agents):
    x = Inches(0.4) + (i * Inches(4.2))
    rrect(s, x, Inches(1.4), Inches(3.95), Inches(5.5), WHITE, accent, Pt(2))
    rect(s, x, Inches(1.4), Inches(3.95), Inches(0.65), accent)
    txbox(s, x + Inches(0.1), Inches(1.43), Inches(3.75), Inches(0.6), f"{icon}  {title}", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    rrect(s, x + Inches(0.15), Inches(2.2), Inches(3.65), Inches(1.0), light, accent, Pt(1))
    txbox(s, x + Inches(0.25), Inches(2.22), Inches(3.45), Inches(0.25), "\U0001F4A1 Simple Analogy:", size=9, bold=True, color=accent, align=PP_ALIGN.LEFT)
    txbox(s, x + Inches(0.25), Inches(2.5), Inches(3.45), Inches(0.65), analogy, size=10, color=TEXT_MEDIUM, align=PP_ALIGN.LEFT)

    bullets(s, x + Inches(0.2), Inches(3.4), Inches(3.55), Inches(2.5), items, size=10, color=TEXT_MEDIUM, spacing=Pt(5))

rrect(s, Inches(2.0), Inches(7.0), Inches(9.33), Inches(0.35), PRIMARY_DARK, None, Pt(0))
txbox(s, Inches(2.0), Inches(7.0), Inches(9.33), Inches(0.35), "Final Decision = Average of all 3 agents' votes (Consensus Approach)", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 7: XAI ENGINE (Visual with icons and examples)
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header_bar(s, "How We Explain Decisions", "Making AI transparent with SHAP, What-If & Stress Testing")
slide_num(s, 7)

xai_items = [
    ("\U0001F4CA", "SHAP Values", ACCENT_GREEN, ACCENT_GREEN_LIGHT,
     "Shows which factors helped or hurt your loan chances",
     ["Green bar (+) = Good for your approval", "Red bar (-) = Hurt your chances", "Longer bar = Bigger impact", "Every feature gets a score"],
     "Example: High Credit Score = +0.15 (helps), Low Income = -0.20 (hurts)"),
    ("\U0001F4A1", "What-If Analysis", ACCENT_AMBER, ACCENT_AMBER_LIGHT,
     "Tells you EXACTLY what to change to get approved",
     ["What if Credit Score was 50 points higher?", "What if Income was 10% more?", "What if Loan Amount was smaller?", "Shows your path to approval"],
     "Example: \"Increase your Credit Score by 30 points to reach 60% approval\""),
    ("\U0001F321", "Sensitivity Heat Map", ACCENT_PURPLE, ACCENT_PURPLE_LIGHT,
     "A color grid showing approval chances across different profiles",
     ["Green zone = High approval chance", "Red zone = Low approval chance", "X-axis: Credit Score", "Y-axis: Income level"],
     "Visual map: See your position and where you need to move"),
    ("\U0001F6E1", "Robustness Score", ACCENT_RED, ACCENT_RED_LIGHT,
     "Tests if the AI decision is stable or easily changed",
     ["Score from 0 to 100", "70+ = Decision is very stable", "40-70 = Somewhat sensitive", "Below 40 = Small changes can flip it"],
     "Example: Score of 85 means the decision is reliable and consistent"),
]

for i, (icon, title, accent, light, subtitle, items, example) in enumerate(xai_items):
    col = i % 2
    row = i // 2
    x = Inches(0.4) + (col * Inches(6.35))
    y = Inches(1.4) + (row * Inches(2.9))
    rrect(s, x, y, Inches(6.1), Inches(2.7), WHITE, accent, Pt(2))
    icon_circle(s, x + Inches(0.2), y + Inches(0.2), Inches(0.65), accent, icon, WHITE, 18)
    txbox(s, x + Inches(1.0), y + Inches(0.15), Inches(4.9), Inches(0.35), title, size=14, bold=True, color=accent)
    txbox(s, x + Inches(1.0), y + Inches(0.5), Inches(4.9), Inches(0.35), subtitle, size=10, color=TEXT_MEDIUM)
    bullets(s, x + Inches(0.3), y + Inches(0.9), Inches(5.5), Inches(1.0), items, size=9, color=TEXT_DARK, spacing=Pt(2))
    rrect(s, x + Inches(0.15), y + Inches(2.1), Inches(5.8), Inches(0.45), light, accent, Pt(1))
    txbox(s, x + Inches(0.25), y + Inches(2.15), Inches(5.6), Inches(0.35), f"\u2192 {example}", size=8, bold=True, color=accent)

# =============================================================================
# SLIDE 8: DASHBOARD SCREENS (Visual mockup descriptions)
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header_bar(s, "The Dashboard", "What the user sees and interacts with")
slide_num(s, 8)

panels = [
    ("\U0001F4DD", "Input Form (Left Side)", PRIMARY, [
        "Enter your financial details",
        "Credit Score slider (300-850)",
        "Choose employment type",
        "Select education level",
        "Enter income & loan amount",
    ]),
    ("\U0001F3AF", "Top Cards (Results)", ACCENT_CYAN, [
        "APPROVED or REJECTED badge",
        "Confidence gauge (0-100%)",
        "Lending Grade: EXL / MOD / OPT",
        "SBI Special Offer if eligible",
    ]),
    ("\U0001F4AC", "AI Summary", ACCENT_GREEN, [
        "Plain English explanation",
        "Typewriter animation effect",
        "Tells you WHY in simple words",
        "Like having an AI advisor talk to you",
    ]),
    ("\U0001F4CA", "SHAP Bar Chart", ACCENT_PURPLE, [
        "Horizontal bars for each feature",
        "Green bars = positive impact",
        "Red bars = negative impact",
        "See at a glance what matters most",
    ]),
    ("\U0001F465", "Agent Jury Panel", ACCENT_AMBER, [
        "Shows all 3 agents' votes",
        "Each agent's confidence level",
        "Progress bars for visual clarity",
        "Approve or Reject per agent",
    ]),
    ("\U0001F3AF", "What-If Sandbox", ACCENT_RED, [
        "Toggle sandbox mode ON",
        "Change any value instantly",
        "See probability update LIVE",
        "Find the optimal profile",
    ]),
]

for i, (icon, title, accent, items) in enumerate(panels):
    col = i % 3
    row = i // 3
    x = Inches(0.4) + (col * Inches(4.2))
    y = Inches(1.4) + (row * Inches(2.9))
    rrect(s, x, y, Inches(3.95), Inches(2.65), WHITE, accent, Pt(2))
    rect(s, x, y, Inches(3.95), Inches(0.5), accent)
    txbox(s, x + Inches(0.1), y + Inches(0.05), Inches(3.75), Inches(0.45), f"{icon}  {title}", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bullets(s, x + Inches(0.15), y + Inches(0.65), Inches(3.65), Inches(1.8), items, size=10, color=TEXT_MEDIUM, spacing=Pt(4))

# =============================================================================
# SLIDE 9: RESULTS (Big numbers, visual achievements)
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header_bar(s, "Results & Impact", "What our framework achieved")
slide_num(s, 9)

big_numbers = [
    ("96.5%", "Accuracy", "3 AI models\nworking together", ACCENT_GREEN),
    ("~42ms", "Speed", "Results in the\nblink of an eye", ACCENT_CYAN),
    ("100%", "Explainable", "Every decision\nfully transparent", PRIMARY),
    ("9", "Input Features", "Complete financial\nprofile analyzed", ACCENT_PURPLE),
]

for i, (value, label, desc, accent) in enumerate(big_numbers):
    x = Inches(0.4) + (i * Inches(3.15))
    rrect(s, x, Inches(1.4), Inches(2.95), Inches(2.2), WHITE, accent, Pt(3))
    txbox(s, x + Inches(0.1), Inches(1.5), Inches(2.75), Inches(0.8), value, size=36, bold=True, color=accent, align=PP_ALIGN.CENTER)
    txbox(s, x + Inches(0.1), Inches(2.3), Inches(2.75), Inches(0.35), label, size=14, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    txbox(s, x + Inches(0.1), Inches(2.7), Inches(2.75), Inches(0.6), desc, size=10, color=TEXT_MEDIUM, align=PP_ALIGN.CENTER)

txbox(s, Inches(0.5), Inches(3.85), Inches(12.33), Inches(0.4), "Key Highlights", size=16, bold=True, color=PRIMARY_DARK, align=PP_ALIGN.CENTER)

highlights = [
    ("\u2705", "3 diverse AI models vote together for higher accuracy than any single model", ACCENT_GREEN),
    ("\U0001F50D", "SHAP values mathematically prove which features drove each decision", ACCENT_CYAN),
    ("\U0001F4A1", "Rejected applicants get clear 'What to change' recommendations", ACCENT_AMBER),
    ("\U0001F4AC", "Complex ML outputs translated to simple English anyone can understand", PRIMARY),
    ("\U0001F4C4", "One-click PDF reports with full audit trail for regulatory compliance", ACCENT_PURPLE),
    ("\U0001F6E1", "Robustness testing ensures predictions are stable and reliable", ACCENT_RED),
]

for i, (icon, text, accent) in enumerate(highlights):
    y = Inches(4.35) + (i * Inches(0.48))
    rrect(s, Inches(0.5), y, Inches(12.33), Inches(0.42), WHITE, BORDER, Pt(1))
    rect(s, Inches(0.5), y, Inches(0.06), Inches(0.42), accent)
    txbox(s, Inches(0.7), y + Inches(0.02), Inches(11.9), Inches(0.38), f"{icon}   {text}", size=11, color=TEXT_DARK)

# =============================================================================
# SLIDE 10: CONCLUSION & THANK YOU
# =============================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)

rect(s, Inches(0), Inches(0), Inches(13.33), Inches(0.12), ACCENT_CYAN)
rect(s, Inches(0), Inches(7.38), Inches(13.33), Inches(0.12), ACCENT_CYAN)

rrect(s, Inches(0.5), Inches(0.5), Inches(6.0), Inches(3.3), PRIMARY, ACCENT_CYAN, Pt(2))
txbox(s, Inches(0.7), Inches(0.6), Inches(5.6), Inches(0.45), "\u2705  What We Built", size=18, bold=True, color=WHITE)

conclusions = [
    "A complete AI system that predicts credit risk with 96.5% accuracy",
    "Every prediction comes with a full, transparent explanation",
    "Interactive dashboard for real-time exploration and testing",
    "Automated PDF reports for compliance and documentation",
    "Proves that AI can be both powerful AND understandable",
]
for i, c in enumerate(conclusions):
    txbox(s, Inches(0.9), Inches(1.2) + (i * Inches(0.42)), Inches(5.4), Inches(0.4), f"\u2713  {c}", size=10, color=ACCENT_CYAN_LIGHT)

rrect(s, Inches(6.9), Inches(0.5), Inches(6.0), Inches(3.3), ACCENT_AMBER, ACCENT_AMBER, Pt(2))
txbox(s, Inches(7.1), Inches(0.6), Inches(5.6), Inches(0.45), "\U0001F680  What's Next", size=18, bold=True, color=PRIMARY_DARK)

futures = [
    "Add more explanation methods (LIME, Anchors)",
    "Include advanced models (XGBoost, LightGBM)",
    "Cloud deployment (AWS / GCP / Azure)",
    "Fairness checks across demographic groups",
    "Real-time model monitoring & drift detection",
]
for i, f in enumerate(futures):
    txbox(s, Inches(7.3), Inches(1.2) + (i * Inches(0.42)), Inches(5.4), Inches(0.4), f"\u2192  {f}", size=10, color=PRIMARY_DARK)

rrect(s, Inches(2.0), Inches(4.1), Inches(9.33), Inches(2.8), PRIMARY_DARK, ACCENT_CYAN, Pt(3))
txbox(s, Inches(2.2), Inches(4.3), Inches(8.93), Inches(0.8), "Thank You!", size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(s, Inches(5.5), Inches(5.15), Inches(2.33), Inches(0.04), ACCENT_CYAN)

txbox(s, Inches(2.2), Inches(5.4), Inches(8.93), Inches(0.4), "Questions & Discussion Welcome", size=16, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

txbox(s, Inches(2.2), Inches(5.9), Inches(8.93), Inches(0.35), "Bibhu Prasad Sahoo  &  Soumyashree", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, Inches(2.2), Inches(6.3), Inches(8.93), Inches(0.3), "8th Semester Research Project", size=11, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

for i, icon in enumerate(["\U0001F9E0", "\U0001F4CA", "\u2699", "\U0001F4AC"]):
    icon_circle(s, Inches(4.1) + Inches(i * 1.3), Inches(6.7), Inches(0.4), ACCENT_CYAN, icon, PRIMARY_DARK, 14)

# =============================================================================
prs.save(OUTPUT_PATH)
print(f"Saved: {OUTPUT_PATH}")
print(f"Size: {os.path.getsize(OUTPUT_PATH) / 1024:.1f} KB")
